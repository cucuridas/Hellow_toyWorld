import win32com.client
from pptx import Presentation
import copy

#path 값으로는 절대위치의 파일명이 들어와야 합니다
#저장 위치 정보에 대해서 어떻게 처리할 것인지....



def solve_excels(_path) :
    #시트가 많은 엑셀 파일을 해제할때
    excel = win32com.client.Dispatch("Excel.Application")
    excel.Visible = True
    count = 1
    wb = excel.Workbooks.Open(_path)
    
    
    while True :
        try :
            print(count)
            sh = wb.Worksheets(count)
            sh.Copy()
            count =count + 1
        except:
            print("해제완료 - 해당내용은 복사만 가능합니다")
            break
        
def solve_excel(_path,number) :
    
    excel = win32com.client.Dispatch("Excel.Application")
    excel.Visible = True
    count = 1
    wb = excel.Workbooks.Open(_path)
    sh = wb.Worksheets(number)
    sh.Copy()
    
    return "해제완료 - 해당내용은 복사만 가능합니다"    
  
#copy 및 복사가 이루어지지 않은 채로 해당 문서를 다른이름으로 저장해버리는 행위(확장자명만 바꾸어 진행)    
def solve_Word(_path) :
    word = win32com.client.Dispatch("Word.Application")
    word.Visible = True

    doc = word.Documents.Open(_path)
    doc.SaveAs('C:\parseData\wordTest.docx_unpaked')
    #doc.Copy()


    print("진행중")


#파일 자체가 ppt로 인식을 해버리는 문제가 있음 -> DRM에서 인식할듯
def solve_PPT(_path):
    powerpt = win32com.client.Dispatch("Powerpoint.Application")
    powerpt.Visible = True
    pt = powerpt.Presentations.Open(_path)
    
    sl = pt.Slides(1)
    sl.Copy()
    #pt.SaveAs('C:\parseData\pptTest.pptx_unpaked123')


    print("개발필요")

def solve_PPT_to_PPTX(_path):
    prs = Presentation(_path)
    cpprs = Presentation()

    # 첫번째 슬라이드 가져오기
    source_slide = prs.slides[0]

    # 새로 추가할 슬라이드 레이아웃 정의 (6 : 빈 페이지)
    slide_layout = cpprs.slide_layouts[0]

    # 새로운 슬라이드 추가
    copy_slide = cpprs.slides.add_slide(slide_layout)

    # 기존 슬라이드(소스)에서 shape를 가져와 새로운 슬라이드에 복사
    for shape in source_slide.shapes:

        el = shape.element

        newel = copy.deepcopy(el)

        copy_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')



#solve_excel('C:\parseData\작업시나리오_GFR-100148_OS 설정 추가_211222.xlsx',17)

solve_PPT_to_PPTX('C:\\parseData\\12_testfor.pptx')